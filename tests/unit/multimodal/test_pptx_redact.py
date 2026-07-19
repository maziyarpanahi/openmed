"""Tests for PPTX slide, table, and speaker-notes redaction."""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pytest

import openmed.multimodal.base as base
from openmed.multimodal import (
    ExtractedDocument,
    extract_pptx,
    map_text_spans_to_pptx_runs,
    redact_document,
    write_redacted_pptx,
)

pptx = pytest.importorskip("pptx")
from pptx.util import Inches  # noqa: E402


@pytest.fixture
def pptx_deps_present(monkeypatch: pytest.MonkeyPatch) -> None:
    """Pretend the complete multimodal extra is installed for dispatch tests."""
    monkeypatch.setattr(base, "_missing_multimodal_dependencies", lambda: [])


@pytest.fixture
def synthetic_pptx(tmp_path: Path) -> Path:
    """Create a synthetic two-slide deck containing no real patient data."""
    presentation = pptx.Presentation()
    blank_layout = presentation.slide_layouts[6]

    first_slide = presentation.slides.add_slide(blank_layout)
    text_box = first_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(8), Inches(1)
    )
    paragraph = text_box.text_frame.paragraphs[0]
    paragraph.add_run().text = "Patient "
    paragraph.add_run().text = "Jane"
    paragraph.add_run().text = " Doe"

    table = first_slide.shapes.add_table(
        1, 1, Inches(0.5), Inches(2), Inches(4), Inches(1)
    ).table
    table_frame = table.cell(0, 0).text_frame
    table_frame.clear()
    table_paragraph = table_frame.paragraphs[0]
    table_paragraph.add_run().text = "MRN "
    table_paragraph.add_run().text = "A123"

    notes_frame = first_slide.notes_slide.notes_text_frame
    assert notes_frame is not None
    notes_frame.clear()
    notes_paragraph = notes_frame.paragraphs[0]
    notes_paragraph.add_run().text = "Call Dr. "
    notes_paragraph.add_run().text = "Alice"
    notes_paragraph.add_run().text = " Smith"

    second_slide = presentation.slides.add_slide(blank_layout)
    follow_up = second_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(8), Inches(1)
    )
    follow_up.text_frame.paragraphs[0].add_run().text = "Follow-up for Bob Stone"
    assert second_slide.has_notes_slide is False

    path = tmp_path / "synthetic_phi.pptx"
    presentation.save(path)
    return path


def _text_parts(path: Path) -> dict[str, list[str]]:
    presentation = pptx.Presentation(path)
    parts: dict[str, list[str]] = {"shapes": [], "tables": [], "notes": []}
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_table:
                parts["tables"].extend(
                    cell.text
                    for row in shape.table.rows
                    for cell in row.cells
                    if cell.text
                )
            elif shape.has_text_frame and shape.text:
                parts["shapes"].append(shape.text)
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None and notes_frame.text:
                parts["notes"].append(notes_frame.text)
    return parts


def _entity(document: ExtractedDocument, text: str, label: str) -> dict[str, Any]:
    start = document.text.index(text)
    return {
        "start": start,
        "end": start + len(text),
        "label": label,
        "confidence": 0.99,
    }


def test_extract_pptx_maps_shape_table_and_notes_offsets_per_slide(
    synthetic_pptx: Path,
) -> None:
    document = extract_pptx(synthetic_pptx)

    first_slide_text = "Patient Jane Doe\nMRN A123\nCall Dr. Alice Smith"
    assert isinstance(document, ExtractedDocument)
    assert document.text == f"{first_slide_text}\nFollow-up for Bob Stone"
    assert document.metadata["format"] == "pptx"
    assert document.metadata["slide_count"] == 2
    assert document.metadata["slide_offsets"] == [
        {
            "slide_index": 0,
            "start": 0,
            "end": len(first_slide_text),
            "length": len(first_slide_text),
        },
        {
            "slide_index": 1,
            "start": len(first_slide_text) + 1,
            "end": len(document.text),
            "length": len("Follow-up for Bob Stone"),
        },
    ]

    for source in document.spans:
        assert document.text_for(source) == document.text[source.start : source.end]
        assert document.location_at(source.start) == source
        assert source.page == source.metadata["slide_index"]

    shape_source = document.location_at(document.text.index("Jane"))
    assert shape_source is not None
    assert document.text_for(shape_source) == "Jane"
    assert shape_source.metadata["block_type"] == "shape"
    assert shape_source.metadata["part"] == "slide"
    assert shape_source.metadata["slide_start"] == len("Patient ")

    table_source = document.location_at(document.text.index("A123"))
    assert table_source is not None
    assert table_source.metadata["block_type"] == "table_cell"
    assert table_source.metadata["row_index"] == 0
    assert table_source.metadata["cell_index"] == 0

    notes_source = document.location_at(document.text.index("Alice"))
    assert notes_source is not None
    assert notes_source.metadata["part"] == "notes"
    assert notes_source.metadata["block_type"] == "speaker_notes"
    assert notes_source.metadata["slide_start"] == first_slide_text.index("Alice")

    follow_up = document.location_at(document.text.index("Follow-up"))
    assert follow_up is not None
    assert follow_up.page == 1
    assert follow_up.metadata["slide_start"] == 0


def test_pptx_redaction_offsets_round_trip_across_runs(
    synthetic_pptx: Path,
) -> None:
    document = extract_pptx(synthetic_pptx)
    entity = _entity(document, "Jane Doe", "PERSON")

    redactions = map_text_spans_to_pptx_runs(document, [entity])

    assert len(redactions) == 1
    redaction = redactions[0]
    assert redaction.replacement == "[PERSON]"
    assert redaction.metadata["slide_indices"] == [0]
    assert redaction.metadata["source_span_count"] == 2
    assert [item.run_start for item in redaction.run_ranges] == [0, 0]
    assert [item.run_end for item in redaction.run_ranges] == [4, 4]
    assert redaction.run_ranges[0].metadata["slide_start"] == len("Patient ")
    assert redaction.run_ranges[-1].metadata["slide_end"] == len("Patient Jane Doe")


def test_write_redacted_pptx_updates_shapes_tables_and_notes(
    synthetic_pptx: Path, tmp_path: Path
) -> None:
    document = extract_pptx(synthetic_pptx)
    entities = [
        _entity(document, "Jane Doe", "PERSON"),
        _entity(document, "A123", "ID_NUM"),
        _entity(document, "Alice Smith", "PERSON"),
    ]
    output = tmp_path / "redacted.pptx"

    assert write_redacted_pptx(synthetic_pptx, output, entities) == output

    parts = _text_parts(output)
    assert "Patient [PERSON]" in parts["shapes"]
    assert "MRN [ID_NUM]" in parts["tables"]
    assert "Call Dr. [PERSON]" in parts["notes"]
    assert "Follow-up for Bob Stone" in parts["shapes"]
    assert "Jane Doe" not in str(parts)
    assert "A123" not in str(parts)
    assert "Alice Smith" not in str(parts)


def test_pptx_handler_is_registered_and_can_write_a_redacted_copy(
    synthetic_pptx: Path,
    tmp_path: Path,
    pptx_deps_present: None,
) -> None:
    output = tmp_path / "dispatched.pptx"

    def detector(text: str, *, lang: str | None = None) -> dict[str, Any]:
        assert lang == "en"
        start = text.index("Alice Smith")
        return {
            "entities": [
                {
                    "start": start,
                    "end": start + len("Alice Smith"),
                    "label": "PERSON",
                    "confidence": 0.98,
                }
            ]
        }

    document = redact_document(
        synthetic_pptx,
        models={"detector": detector},
        lang="en",
        policy={"output_path": output},
    )

    assert ".pptx" in base._HANDLERS
    assert document.metadata["detected_span_count"] == 1
    assert document.metadata["redacted_pptx_path"] == str(output)
    redaction = document.metadata["pptx_redactions"][0]
    assert redaction["label"] == "PERSON"
    assert redaction["run_ranges"][0]["metadata"]["part"] == "notes"
    assert "Alice Smith" not in str(redaction)
    assert "Call Dr. [PERSON]" in _text_parts(output)["notes"]
